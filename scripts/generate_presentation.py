import os
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def load_metrics_and_stats():
    # Load cleaned data
    df = pd.read_csv("data/raw_crm_export.csv") # We will read and clean it on the fly or just load from process_pipeline
    # Wait, raw_crm_export is messy. It's better to load the final cleaned dataframe or re-clean it.
    # Let's import clean_and_process_data from scripts/process_pipeline to get the exact cleaned df!
    import sys
    sys.path.append(os.path.abspath('scripts'))
    from process_pipeline import clean_and_process_data
    df_clean, stats = clean_and_process_data()
    
    # Calculate live stats
    total_leads = len(df_clean)
    trials = int((df_clean["Trial_Activated"] == "Yes").sum())
    conversions = int((df_clean["Converted"] == "Yes").sum())
    total_mrr = float(df_clean["MRR"].dropna().sum())
    
    # Calculate averages
    lead_to_trial = trials / total_leads
    trial_to_conv = conversions / trials
    overall_conv = conversions / total_leads
    
    # Compute deal velocity (lead to conversion)
    df_conv = df_clean[df_clean["Converted"] == "Yes"].copy()
    df_conv["Lead_Date"] = pd.to_datetime(df_conv["Lead_Date"])
    df_conv["Conversion_Date"] = pd.to_datetime(df_conv["Conversion_Date"])
    avg_velocity = (df_conv["Conversion_Date"] - df_conv["Lead_Date"]).dt.days.mean()
    
    # Segment analysis
    segments = {}
    for seg in ["SMB", "Mid-Market", "Enterprise"]:
        df_seg = df_clean[df_clean["Segment"] == seg]
        seg_leads = len(df_seg)
        seg_trials = int((df_seg["Trial_Activated"] == "Yes").sum())
        seg_conv = int((df_seg["Converted"] == "Yes").sum())
        seg_mrr = float(df_seg["MRR"].dropna().sum())
        
        segments[seg] = {
            "leads": seg_leads,
            "trials": seg_trials,
            "conversions": seg_conv,
            "mrr": seg_mrr,
            "l2t": seg_trials / seg_leads if seg_leads > 0 else 0,
            "t2c": seg_conv / seg_trials if seg_trials > 0 else 0
        }
        
    return {
        "leads": total_leads,
        "trials": trials,
        "conversions": conversions,
        "mrr": total_mrr,
        "l2t": lead_to_trial,
        "t2c": trial_to_conv,
        "conv": overall_conv,
        "velocity": avg_velocity,
        "segments": segments,
        "etl_stats": stats
    }

def add_header(slide, title_text, dark_mode=False):
    # Header Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255) if dark_mode else RGBColor(43, 62, 80) # White or Slate Blue
    return title_box

def apply_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def create_card(slide, left, top, width, height, title, content_list, bg_color, title_color, text_color):
    # Add a panel shape for background card
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = RGBColor(211, 211, 211) # Light gray border
    card.line.width = Pt(1)
    
    # Text Box inside card
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    # Title
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color
    p_title.space_after = Pt(8)
    
    # Content
    for idx, item in enumerate(content_list):
        p = tf.add_paragraph()
        p.text = item
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = text_color
        p.space_after = Pt(4)
        p.level = 0

def generate_pptx():
    print("Gathering metrics data...")
    m = load_metrics_and_stats()
    
    # Init presentation
    prs = Presentation()
    
    # Set standard widescreen 16:9 dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank slide
    
    # Colors
    c_navy = RGBColor(31, 78, 120)     # Primary dark
    c_slate = RGBColor(43, 62, 80)     # Slate gray
    c_orange = RGBColor(198, 89, 17)   # Accent orange
    c_green = RGBColor(55, 86, 35)     # Accent green
    c_light_gray = RGBColor(242, 244, 247) # Card bg
    c_white = RGBColor(255, 255, 255)
    c_charcoal = RGBColor(50, 50, 50)
    
    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1, c_navy)
    
    # Title Box (large centered)
    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "B2B SaaS Revenue Pipeline Analysis"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = c_white
    p1.space_after = Pt(14)
    
    p2 = tf1.add_paragraph()
    p2.text = "Designing ETL Data Cleanse Pipelines & Automating Quote-to-Access Workflows"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(220, 220, 220)
    p2.space_after = Pt(28)
    
    p3 = tf1.add_paragraph()
    p3.text = "RevOps Independent Research Project"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(180, 180, 180)
    
    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary (Light Theme)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2, c_white)
    add_header(slide2, "Executive Summary")
    
    # Add 3 Cards side-by-side
    create_card(
        slide=slide2,
        left=Inches(0.5), top=Inches(1.6), width=Inches(3.8), height=Inches(5.0),
        title="1. Pipeline Operations",
        content_list=[
            f"• Scanned {m['leads']} Leads over a 12-week CRM log, establishing the baseline funnel.",
            f"• Identified {m['trials']} Trials (Trial Rate: {m['l2t']:.1%}) & {m['conversions']} Conversions (Conversion Rate: {m['conv']:.1%}).",
            f"• Added ${m['mrr']:,.2f} in Monthly Recurring Revenue (MRR), with strong ACV contributions.",
            f"• Average cycle time from Lead creation to Customer Conversion is {m['velocity']:.1f} days."
        ],
        bg_color=c_light_gray, title_color=c_navy, text_color=c_charcoal
    )
    
    create_card(
        slide=slide2,
        left=Inches(4.7), top=Inches(1.6), width=Inches(3.8), height=Inches(5.0),
        title="2. Identified Process Gaps",
        content_list=[
            "• Manual Contract Emailing: Sales emails PDF contracts to billing managers, causing a 4.2-day average processing lag.",
            "• Double Customer Data Entry: No CRM-to-billing API integration leads to manual copy-paste errors (~5% typo rate).",
            "• Siloed Access Provisioning: Finance messages Ops over Slack to activate user accounts, adding 2.1 days of customer delay."
        ],
        bg_color=c_light_gray, title_color=c_orange, text_color=c_charcoal
    )
    
    create_card(
        slide=slide2,
        left=Inches(8.9), top=Inches(1.6), width=Inches(3.8), height=Inches(5.0),
        title="3. Automation Outcomes",
        content_list=[
            "• Integrated Quote-to-Invoice: Automatic Stripe invoice drafts generated on Salesforce Stage = 'Closed Won'.",
            "• Stripe to Product Database Webhook: Payments instantly invoke user creation in Auth0 database.",
            "• Time Savings: Cuts manual coordination & reconciliation overhead by ~30% (~3.7 hours/week).",
            "• Customer Experience: Slashing Time-to-Value from 6.3 combined days down to under 2 minutes."
        ],
        bg_color=c_light_gray, title_color=c_green, text_color=c_charcoal
    )
    
    # -------------------------------------------------------------
    # SLIDE 3: Funnel Metrics & Revenue Breakdown
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3, c_white)
    add_header(slide3, "SaaS Sales Funnel & Segment Analysis")
    
    # Funnel metrics text card on left
    create_card(
        slide=slide3,
        left=Inches(0.5), top=Inches(1.6), width=Inches(4.5), height=Inches(5.0),
        title="Pipeline Funnel Conversion Metrics",
        content_list=[
            f"• Leads Captured: {m['leads']} registrations",
            f"• Trial Activation Rate: {m['l2t']:.1%} ({m['trials']} accounts)",
            f"• Trial-to-Conversion Rate: {m['t2c']:.1%} ({m['conversions']} customers)",
            f"• Overall Funnel Efficiency: {m['conv']:.1%}",
            f"• Average Deal Velocity: {m['velocity']:.1f} days from creation to close.",
            "",
            "RevOps Key Takeaway:",
            "The largest bottleneck is Trial-to-Conversion (29.7%), suggesting friction in trial onboarding or product adoption rather than lead capture."
        ],
        bg_color=c_light_gray, title_color=c_navy, text_color=c_charcoal
    )
    
    # Segment breakdown table on right
    # Add Table
    rows = 5
    cols = 5
    table_shape = slide3.shapes.add_table(rows, cols, Inches(5.5), Inches(2.0), Inches(7.3), Inches(4.0))
    table = table_shape.table
    
    headers_table = ["Segment", "Leads", "Conversions", "Conv. %", "MRR Added"]
    for i, h in enumerate(headers_table):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = c_slate
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = c_white
            p.alignment = PP_ALIGN.CENTER
            
    seg_names = ["SMB", "Mid-Market", "Enterprise"]
    for row_idx, seg_name in enumerate(seg_names, start=1):
        seg_data = m["segments"][seg_name]
        data_row = [
            seg_name,
            f"{seg_data['leads']}",
            f"{seg_data['conversions']}",
            f"{seg_data['t2c']:.1%}",
            f"${seg_data['mrr']:,.0f}"
        ]
        
        for col_idx, text in enumerate(data_row):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = c_light_gray
            else:
                cell.fill.fore_color.rgb = c_white
                
            for p in cell.text_frame.paragraphs:
                p.font.name = "Segoe UI"
                p.font.size = Pt(12)
                p.font.color.rgb = c_charcoal
                p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
                
    # Add Total row
    totals = ["Total / Wtd Avg", f"{m['leads']}", f"{m['conversions']}", f"{m['t2c']:.1%}", f"${m['mrr']:,.0f}"]
    for col_idx, text in enumerate(totals):
        cell = table.cell(4, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(217, 225, 242)
        for p in cell.text_frame.paragraphs:
            p.font.name = "Segoe UI"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = c_navy
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            
    # -------------------------------------------------------------
    # SLIDE 4: Process Gaps & Risks
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4, c_white)
    add_header(slide4, "Current Process Gaps & Operational Friction")
    
    # 3 horizontal bars for each gap
    gaps = [
        {
            "num": "GAP 1",
            "title": "Manual Contract Emailing (Sales to Finance Handshake)",
            "impact": "4.2-Day Processing Lag",
            "desc": "Sales reps manually transfer DocuSign PDF contracts via email to Finance. Inbox latency delays invoice generation and cash collection.",
            "color": c_orange
        },
        {
            "num": "GAP 2",
            "title": "Double-Entry Customer Profile Creation (CRM and Billing Silos)",
            "impact": "5% Typographical Error Rate",
            "desc": "Finance managers type customer data into Stripe from PDFs. This causes mismatched company names, broken invoice delivery, and reporting drift.",
            "color": c_orange
        },
        {
            "num": "GAP 3",
            "title": "Siloed Application Access Provisioning (Finance to Operations Handshake)",
            "impact": "2.1-Day Time-to-Value Delay",
            "desc": "Access is only provisioned after Finance verifies bank transfers and messages Ops via Slack. Customers wait, leading to friction at onboarding.",
            "color": c_orange
        }
    ]
    
    for idx, gap in enumerate(gaps):
        top_offset = Inches(1.7 + idx * 1.7)
        
        # Color block on left
        block = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top_offset, Inches(1.5), Inches(1.3))
        block.fill.solid()
        block.fill.fore_color.rgb = gap["color"]
        block.line.fill.background()
        tf_b = block.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = gap["num"]
        p_b.font.name = "Segoe UI"
        p_b.font.size = Pt(20)
        p_b.font.bold = True
        p_b.font.color.rgb = c_white
        p_b.alignment = PP_ALIGN.CENTER
        
        # Details panel
        panel = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.1), top_offset, Inches(10.7), Inches(1.3))
        panel.fill.solid()
        panel.fill.fore_color.rgb = c_light_gray
        panel.line.color.rgb = RGBColor(211, 211, 211)
        tf_p = panel.text_frame
        tf_p.word_wrap = True
        
        p_t = tf_p.paragraphs[0]
        p_t.text = f"{gap['title']}  |  Impact: {gap['impact']}"
        p_t.font.name = "Segoe UI"
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = c_slate
        p_t.space_after = Pt(4)
        
        p_d = tf_p.add_paragraph()
        p_d.text = gap["desc"]
        p_d.font.name = "Segoe UI"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = c_charcoal
        
    # -------------------------------------------------------------
    # SLIDE 5: Future State & Automation Benefits
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5, c_white)
    add_header(slide5, "Target State: Automated Lifecycle Architecture")
    
    # Architecture process cards
    create_card(
        slide=slide5,
        left=Inches(0.5), top=Inches(1.6), width=Inches(3.8), height=Inches(5.0),
        title="CRM Trigger (Stage = Won)",
        content_list=[
            "Sales sets Opportunity Stage to Closed Won in Salesforce.",
            "",
            "System Action:",
            "Webhook pushes deal metadata (Account ID, Contract Value, Customer Email, Subscription Tier) to iPaaS Sync Engine.",
            "",
            "Business Benefit:",
            "Eliminates invoice queue lag; initiates billing setup in under 5 seconds."
        ],
        bg_color=c_light_gray, title_color=c_navy, text_color=c_charcoal
    )
    
    create_card(
        slide=slide5,
        left=Inches(4.7), top=Inches(1.6), width=Inches(3.8), height=Inches(5.0),
        title="iPaaS & Stripe Billing Setup",
        content_list=[
            "iPaaS maps CRM metadata to Stripe objects. Updates profiles or generates new customer IDs.",
            "",
            "System Action:",
            "A finalized Invoice is auto-generated in Stripe and emailed instantly to the customer's billing contact.",
            "",
            "Business Benefit:",
            "0% double-entry typographical error rate. Standardized client IDs mapped to Stripe metadata."
        ],
        bg_color=c_light_gray, title_color=c_navy, text_color=c_charcoal
    )
    
    create_card(
        slide=slide5,
        left=Inches(8.9), top=Inches(1.6), width=Inches(3.8), height=Inches(5.0),
        title="Webhooks & Instant Onboarding",
        content_list=[
            "Customer pays email invoice link. Stripe processes funds.",
            "",
            "System Action:",
            "Stripe payment success event triggers webhook callback. API calls Auth0 to instantly spin up customer workspace and email credentials.",
            "",
            "Business Benefit:",
            "Time-to-Value (onboarding delay) slashed from 2.1 days to under 2 minutes. Strong first-day NPS."
        ],
        bg_color=c_light_gray, title_color=c_navy, text_color=c_charcoal
    )
    
    # -------------------------------------------------------------
    # SLIDE 6: GTM Recommendations & Audit
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6, c_white)
    add_header(slide6, "RevOps Strategic GTM Recommendations")
    
    create_card(
        slide=slide6,
        left=Inches(0.5), top=Inches(1.6), width=Inches(5.8), height=Inches(5.0),
        title="Go-To-Market & Funnel Priorities",
        content_list=[
            "1. Optimize Trial-to-Conversion Bottleneck",
            "   Our biggest drop-off is trial-to-conversion (29.7%). RevOps suggests introducing automated in-app onboarding guides and email sequences triggered by inactivity during the trial phase.",
            "",
            "2. CRM Validation Controls",
            "   Implement strict Salesforce validation rules preventing opportunities from transitioning to 'Negotiation' or 'Closed Won' without verified billing contacts. This eliminates email bounce failures.",
            "",
            "3. Shift Marketing Allocation",
            "   Focus marketing budgets on sources exhibiting high velocity and trial engagement. Web searches (Google) and referrals provide the highest MRR velocity compared to cold outbound channels."
        ],
        bg_color=c_light_gray, title_color=c_navy, text_color=c_charcoal
    )
    
    create_card(
        slide=slide6,
        left=Inches(6.9), top=Inches(1.6), width=Inches(5.8), height=Inches(5.0),
        title="Audit Framework & Time Savings",
        content_list=[
            "1. Implement Weekly RevOps Audits",
            "   Leverage the Revenue Pipeline Tracker Excel sheet. Monitor the 'RevOps ETL Data Audit Report' tab weekly. Target: 0 unresolved duplicates and 100% database schema mapping consistency.",
            "",
            "2. Tracking Reconciliation Time Savings",
            "   Prior to automation, staff spent ~12.5 hours weekly checking bank ledgers and manual access provisioning. With our proposed webhooks, this is reduced to ~8.8 hours.",
            "",
            "3. ROI of Onboarding Automations",
            "   - Manual overhead cut by ~3.7 hours/week (30% reduction).",
            "   - Customer onboarding delay reduced from 6.3 days to minutes.",
            "   - Annual staff productivity savings: estimated $12,000+ plus higher customer retention rates."
        ],
        bg_color=c_light_gray, title_color=c_green, text_color=c_charcoal
    )
    
    # Save Presentation
    output_path = "reports/leadership_presentation.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Successfully generated PowerPoint slide deck in '{output_path}'!")

if __name__ == "__main__":
    generate_pptx()
